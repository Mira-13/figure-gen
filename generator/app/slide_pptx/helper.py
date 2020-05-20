import os
import json
import numpy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from . import calculate


def add_img_to_slide(slide, path, width_inch, pos_top, pos_left):
    slide.shapes.add_picture(path, Inches(pos_left), Inches(pos_top), width=Inches(width_inch))

def add_frame_on_top(slide, pos_top, pos_left, width_inch, height_inch, color, thickness_pt):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(pos_left), Inches(pos_top), Inches(width_inch), Inches(height_inch)
    )

    shape.line.color.rgb = RGBColor(color[0], color[1], color[2])
    shape.line.width = Pt(thickness_pt)
    shape.fill.background()
    shape.shadow.inherit = False

def has_frame(element):
    try:
        frame = element['frame']
    except:
        return False
    return True

def place_images_and_frames(slide, data, factor, offset_width_mm):
    '''
    Reads module data and puts images on the slide. 
    args:
        factor; is slide_width / figure_width, which is used to position elements on the slide accordingly
    '''
    width_inch, height_inch = calculate.img_size_inches(data, factor)

    rowIndex = 1
    for row in data['elements_content']:
        colIndex = 1
        if rowIndex <= data['num_rows']:
            for element in row:
                if colIndex <= data['num_columns']:
                    pos_top, pos_left = calculate.img_pos_for_slide(data, colIndex, rowIndex, factor, offset_left_mm=offset_width_mm)
                    add_img_to_slide(slide, element['filename'], width_inch, pos_top, pos_left)
                    #print(element['filename'], pos_top, pos_left)
                    if has_frame(element):
                        add_frame_on_top(slide, pos_top, pos_left, width_inch, height_inch, color=element['frame']['color'], thickness_pt=element['frame']['line_width'])
                    colIndex += 1
            rowIndex += 1

def add_text(slide, content, rotation, fontsize, factor, pos_top, pos_left):
    pass # TODO

def place_titles(slide, data, factor):
    pass # TODO

def place_row_titles(slide, data, factor):
    pass # Todo

def place_col_titles(slide, data, factor):
    pass #Todo





