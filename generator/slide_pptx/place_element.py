import os
import json
import numpy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from . import calculate


def add_image(slide, path, width_inch, pos_top, pos_left):
    slide.shapes.add_picture(path, Inches(pos_left), Inches(pos_top), width=Inches(width_inch))

def create_rectangle(slide, pos_top, pos_left, width, height):
    '''
    position and size in inches.
    Default: no shadow
    '''
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(pos_left), Inches(pos_top), Inches(width), Inches(height))
    shape.shadow.inherit = False
    return shape

def apply_fill_color(shape, color):
    if color is None:
        shape.fill.background() # = no fill color
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(color[0], color[1], color[2])
    shape.line.fill.background()

def add_frame_on_top(slide, pos_top, pos_left, width_inch, height_inch, color, thickness_pt):
    '''
    Frames are rectangles placed on top of an image element. 
    A frame has: 
        - no background color, 
        - no text
    '''
    shape = create_rectangle(slide, pos_top, pos_left, width_inch, height_inch)

    shape.line.color.rgb = RGBColor(color[0], color[1], color[2])
    shape.line.width = Pt(thickness_pt)
    shape.fill.background()

def has_frame(element):
    try:
        frame = element['frame']
    except:
        return False
    return True

def images_and_frames(slide, data, factor, offset_width_mm):
    '''
    Reads module data and puts images on the slide. 
    args:
        factor; is slide_width / figure_width, which is used to position elements on the slide accordingly

    TODO: apply offsets to make space for text
    '''
    width_inch, height_inch = calculate.img_size_inches(data, factor)

    rowIndex = 1
    for row in data['elements_content']:
        colIndex = 1
        if rowIndex <= data['num_rows']:
            for element in row:
                if colIndex <= data['num_columns']:
                    pos_top, pos_left = calculate.img_pos_for_slide(data, colIndex, rowIndex, factor, offset_left_mm=offset_width_mm)
                    add_image(slide, element['filename'], width_inch, pos_top, pos_left)
                    #print(element['filename'], pos_top, pos_left)
                    if has_frame(element):
                        add_frame_on_top(slide, pos_top, pos_left, width_inch, height_inch, color=element['frame']['color'], thickness_pt=element['frame']['line_width'])
                    colIndex += 1
            rowIndex += 1

def apply_text_properties(shape, text, fontsize, txt_color):
    '''
    Insert text into shape, set fontsize, and applies text color.
    Default alignment is 'center' and all margins are set to 0
    '''
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    text_frame.margin_left = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.margin_right = 0

    run = p.add_run()
    run.text = text

    font = run.font
    font.color.rgb = RGBColor(txt_color[0], txt_color[1], txt_color[2])
    font.size = Pt(fontsize)

def add_text(slide, pos_top, pos_left, width_inch, height_inch, text, txt_rotation, fontsize_pt, txt_color=[0,0,0], bg_color=None):
    '''
    Add a rectangle which includes:
        - text properties: text content, rotation (only 0° / +- 90°!), text color
        - box properties: background color (but no frame)

    Careful: pptx package does not support text rotation, yet.
    However, we support the rotation (+-) 90° by doing some workarounds:
    - we switch height and width of the rectangle, and 
    - rotate the whole shape afterwards
    --> the shape, including it's content (text) is therefore rotated.
    Other rotation values are not supported, except for 0°. 
    '''
    if txt_rotation == 90.0 or txt_rotation == -90.0:
        shape = create_rectangle(slide, pos_top + height_inch / 2. - width_inch / 2, pos_left - height_inch / 2 + width_inch / 2, height_inch, width_inch)
        shape.rotation = txt_rotation
    else:
        shape = create_rectangle(slide, pos_top, pos_left, width_inch, height_inch)
    apply_fill_color(shape, bg_color)
    apply_text_properties(shape, text, fontsize_pt, txt_color)


def titles(slide, data, factor, offset_left_mm, offset_top_mm=0.0):
    for direction in ['north', 'east', 'south', 'west']:
        position, size = calculate.titles_pos_for_slide(data, direction, factor, offset_left_mm, offset_top_mm)
        if size[0] != 0.0 and size[1] != 0.0:
            title = data['titles'][direction]
            add_text(slide, position[0], position[1], size[0], size[1], title['content'], title['rotation'], title['fontsize'], title['text_color'], title['background_color'])
            
def row_titles(slide, data, factor, offset_left_mm, offset_top_mm=0.0):
    for direction in ['east', 'west']:
        if calculate.padding_of(data['row_titles'], direction)[0] != 0.0:
            bg_color_porperties = data['row_titles'][direction]['background_colors']
            if bg_color_porperties is None:
                single_bg_color = bg_color_porperties # None
            elif not isinstance(bg_color_porperties[0], list): # = if first element is number(int or float) and not list
                single_bg_color = bg_color_porperties
            else:
                single_bg_color = None

            for cur_row in range(1, data['num_rows']+1):
                position, size = calculate.row_titles_pos_for_slide(data, cur_row, direction, factor, offset_left_mm, offset_top_mm)
                r_title = data['row_titles'][direction]
                if single_bg_color is None and bg_color_porperties is not None:
                    bg_color = r_title['background_colors'][cur_row-1]
                else:
                    bg_color = single_bg_color
                add_text(slide, position[0], position[1], size[0], size[1], r_title['content'][cur_row-1], r_title['rotation'], r_title['fontsize'], r_title['text_color'], bg_color)

def col_titles(slide, data, factor, offset_left_mm, offset_top_mm=0.0):
    for direction in ['north', 'south']:
        if calculate.padding_of(data['column_titles'], direction)[0] != 0.0:
            bg_color_porperties = data['column_titles'][direction]['background_colors']
            if bg_color_porperties is None:
                single_bg_color = bg_color_porperties # None
            elif not isinstance(bg_color_porperties[0], list): # = if first element is number(int or float) and not list
                single_bg_color = bg_color_porperties
            else:
                single_bg_color = None

            for cur_row in range(1, data['num_columns']+1):
                position, size = calculate.column_titles_pos_for_slide(data, cur_row, direction, factor, offset_left_mm, offset_top_mm)
                r_title = data['column_titles'][direction]
                if single_bg_color is None and bg_color_porperties is not None:
                    bg_color = r_title['background_colors'][cur_row-1]
                else:
                    bg_color = single_bg_color
                add_text(slide, position[0], position[1], size[0], size[1], r_title['content'][cur_row-1], r_title['rotation'], r_title['fontsize'], r_title['text_color'], bg_color)





