import os
import json
from ..mplot import make_plot
from pptx import Presentation
from pptx.util import Inches
from . import place_element, calculate

'''
in PPTX format we 
 - ignore background colors
 - ignore element captions (north/east/south/west content of each img) as we didn't even use them once before
 - do not support 'dashed' frames - if a frame is 'dashed' the frame in pptx will be normal (but still has a frame)
 - only support text rotation by 0° and +-90°
''' 

class Error(Exception):
    def __init__(self, message):
        self.message = message
    
def generate(module_data, to_path, index, delete_gen_files=True):
    return module_data

def combine(data, to_path, delete_gen_files=True):
    # calculate correct width scaling so that the figure fills out the slide
    sum_total_width_mm = 0
    for d in data:
        sum_total_width_mm += d['total_width']

    #create slide
    prs = Presentation()
    if True:
        figure_height = data[0]['total_height']
        if figure_height < 25.4: # mm
            figure_height = 25.4
            print("Warning: pptx computed height is less than the minimum of 2,54cm. The slide heights will be set on 2,54cm.")
        prs.slide_height = Inches(calculate.mm_to_inch(figure_height)) 
        
        prs.slide_width = Inches(calculate.mm_to_inch(sum_total_width_mm))
        width_scaling = 1
    else:
        prs.slide_height = Inches(9) 
        prs.slide_width = Inches(16)
        width_scaling = 16 / calculate.mm_to_inch(sum_total_width_mm)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)   

    # generate content
    cur_width_mm = 0
    idx = 0
    for d in data:
        if d['type'] == 'plot':
            filename = os.path.join(to_path, f"plot{idx}.png")
            make_plot.generate(d, to_path, f"plot{idx}.png")
            place_element.add_image(slide, filename, calculate.mm_to_inch(d['total_width']), 0, calculate.mm_to_inch(cur_width_mm))
            idx += 1
        else:
            place_element.images_and_frames_and_labels(slide, d, width_scaling, cur_width_mm)
            place_element.titles(slide, d, width_scaling, cur_width_mm)
            place_element.row_titles(slide, d, width_scaling, cur_width_mm)
            place_element.col_titles(slide, d, width_scaling, cur_width_mm)
        cur_width_mm += d['total_width']

    # save
    path_file = os.path.join(to_path, 'gen_figure.pptx')
    prs.save(path_file)