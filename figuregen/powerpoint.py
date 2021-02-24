from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
import importlib.resources as pkg_resources
from concurrent.futures import ThreadPoolExecutor, Future
from threading import Lock
from .backend import *

class PptxBackend(Backend):
    ''' A very basic .pptx generator that only roughly matches results of other backends.

    in PPTX format we have the following limitations due to our dependency python-pptx:
    - ignore background colors
    - ignore vertical alignment and padding of titles
    - do not support 'dashed' frames - if a frame is 'dashed' the frame in pptx will be normal (but still has a frame)
    - only support text rotation by 0° and +-90°
    '''
    def __init__(self):
        self._thread_pool = ThreadPoolExecutor()
        self._slide_mutex = Lock()

    def assemble_grid(self, components: List[Component], output_dir: str):
        return components

    def combine_grids(self, data: List[List[Component]], idx: int, bounds: Bounds) -> List[Component]:
        flat = []
        for row in data:
            flat.extend(row)
        return flat

    def _add_image(self, c: Component, slide):
        # Write image to temp folder
        with tempfile.TemporaryDirectory() as tmpdir:
            fname = c.data.make_raster(c.bounds.width, c.bounds.height, os.path.join(tmpdir, "image"))
            self._slide_mutex.acquire()
            shape = slide.shapes.add_picture(fname, Mm(c.bounds.left), Mm(c.bounds.top),
                width=Mm(c.bounds.width))
            shape.shadow.inherit = False
            self._slide_mutex.release()

        if c.has_frame:
            self._slide_mutex.acquire()
            offset = Pt(c.frame_linewidth) / 2
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(c.bounds.left) + offset,
                Mm(c.bounds.top) + offset, Mm(c.bounds.width) - Pt(c.frame_linewidth),
                Mm(c.bounds.height) - Pt(c.frame_linewidth))
            shape.shadow.inherit = False
            shape.line.color.rgb = RGBColor(c.frame_color[0], c.frame_color[1], c.frame_color[2])
            shape.line.width = Pt(c.frame_linewidth)
            # shape.line.join_type = 'Miter' # Removes rounded edges, but is not supported, yet (sadly)
            shape.fill.background()
            self._slide_mutex.release()

    def combine_rows(self, data: List[Component], bounds: Bounds):
        # We load a template from a file to have some nicer line styles etc by default
        # (they cannot currently be specified via python-pptx)
        with tempfile.TemporaryDirectory() as tmpdir:
            themedata = pkg_resources.read_binary(__package__, "theme.pptx")
            p = os.path.join(tmpdir, "theme.pptx")
            with open(p, "wb") as f:
                f.write(themedata)
            prs = Presentation(p)

        # Create a single slide presentation with a blank slide
        # prs = Presentation()
        prs.slide_height = Mm(bounds.height)
        prs.slide_width = Mm(bounds.width)
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add all our elements to the slide
        flat = []
        for row in data:
            flat.extend(row)

        # Generate all images in parallel
        futures = []
        for c in flat:
            if isinstance(c, ImageComponent):
                futures.append(self._thread_pool.submit(self._add_image, c, slide))
        for f in futures:
            f.result()

        # Add everything else afterwards, to ensure proper z-order
        for c in flat:
            if isinstance(c, TextComponent):
                if c.rotation == 90.0 or c.rotation == -90.0:
                    # The shape is rotated about its center. We want a rotation about the top left corner instead.
                    # Since we only allow 90° rotations, we can correct for that with a simple translation
                    pos_top = c.bounds.top + c.bounds.height / 2. - c.bounds.width / 2.
                    pos_left = c.bounds.left - c.bounds.height / 2. + c.bounds.width / 2.

                    # swap height and width
                    height, width = c.bounds.width, c.bounds.height

                    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(pos_left), Mm(pos_top),
                        Mm(width), Mm(height))
                    # tikz rotation is counter-clockwise, pptx clockwise (we switch in pptx)
                    shape.rotation = -c.rotation
                else:
                    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(c.bounds.left), Mm(c.bounds.top),
                        Mm(c.bounds.width), Mm(c.bounds.height))

                shape.shadow.inherit = False

                # Background color
                if c.background_color is not None:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(c.background_color[0], c.background_color[1],
                        c.background_color[2])
                else:
                    shape.fill.background()
                shape.line.fill.background()

                # Text properties
                text_frame = shape.text_frame
                p = text_frame.paragraphs[0]
                p.alignment = {
                    "center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT
                }[c.horizontal_alignment]

                text_frame.margin_top = 0
                text_frame.margin_bottom = 0

                if c.horizontal_alignment == 'right':
                    text_frame.margin_right = Mm(c.padding.width_mm)
                    text_frame.margin_left = 0
                else:
                    text_frame.margin_right = 0
                    text_frame.margin_left = Mm(c.padding.width_mm)

                run = p.add_run()
                run.text = c.content.replace("\\\\", "\n")
                run.font.color.rgb = RGBColor(c.color[0], c.color[1], c.color[2])
                run.font.size = Pt(c.fontsize)

            if isinstance(c, RectangleComponent):
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(c.bounds.left), Mm(c.bounds.top),
                        Mm(c.bounds.width), Mm(c.bounds.height))
                shape.shadow.inherit = False
                shape.line.color.rgb = RGBColor(c.color[0], c.color[1], c.color[2])
                shape.line.width = Pt(c.linewidth)
                # shape.line.join_type = 'Miter' # Removes rounded edges, but is not supported, yet (sadly)
                shape.fill.background()

            if isinstance(c, LineComponent):
                shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Mm(c.from_x), Mm(c.from_y),
                    Mm(c.to_x), Mm(c.to_y))
                shape.shadow.inherit = False
                shape.line.color.rgb = RGBColor(c.color[0], c.color[1], c.color[2])
                shape.line.width = Pt(c.linewidth)

        return prs

    def write_to_file(self, data, filename: str):
        data.save(filename)