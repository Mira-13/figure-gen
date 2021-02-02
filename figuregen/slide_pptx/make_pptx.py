import os
import threading
from pptx import Presentation
from pptx.util import Inches
from . import place_element, calculate
from ..element_data import *

'''
in PPTX format we
 - ignore background colors
 - do not support 'dashed' frames - if a frame is 'dashed' the frame in pptx will be normal (but still has a frame)
 - only support text rotation by 0° and +-90° (this is a limitation of python-pptx)
'''

class GridError(Exception):
    def __init__(self, row, col, message):
        self.message = f"Error in row {row}, column {col}: {message}"

def _export_image(module, figure_idx, module_idx, path, row, col):
    elem = module["elements_content"][row][col]
    width, height = module['element_config']['img_width'], module['element_config']['img_height']

    assert isinstance(elem["image"], ElementData), "Element is of the wrong type."

    prefix = f'img-{row+1}-{col+1}-{figure_idx+1}-{module_idx+1}'
    prefix = os.path.join(path, prefix)

    elem["image"] = elem["image"].make_raster(width, height, prefix)

def export_images(module, figure_idx, module_idx, path):
    threads = []
    for row in range(module["num_rows"]):
        for col in range(module["num_columns"]):
            t = threading.Thread(target=_export_image, args=(module, figure_idx, module_idx, path, row, col))
            t.start()
            threads.append(t)
    for t in threads:
        t.join()

def generate(module_data, figure_idx, module_idx, temp_folder, delete_gen_files=True, tex_packages=[]):
    export_images(module_data, figure_idx, module_idx, path=temp_folder)
    return module_data

def place_modules(data, to_path, slide):
    offset_top_mm = 0
    for fig_idx in range(len(data)):
        cur_width_mm = 0
        for d in data[fig_idx]:
            place_element.images_and_frames_and_labels(slide, d, 1, cur_width_mm, offset_top_mm)
            place_element.titles(slide, d, 1, cur_width_mm, offset_top_mm)
            place_element.row_titles(slide, d, 1, cur_width_mm, offset_top_mm)
            place_element.col_titles(slide, d, 1, cur_width_mm, offset_top_mm)
            place_element.south_captions(slide, d, 1, cur_width_mm, offset_top_mm)
            cur_width_mm += d['total_width']
        offset_top_mm += data[fig_idx][0]['total_height']

def combine(data, filename, temp_folder, delete_gen_files=True):
    to_path = os.path.dirname(filename)

    # calculate correct width scaling so that the figure fills out the slide
    sum_total_width_mm = 0
    for d in data[0]:
        sum_total_width_mm += d['total_width']

    figure_height = 0.
    for d in data:
        figure_height += d[0]['total_height']

    if figure_height < 25.4: # mm
        figure_height = 25.4
        print("Warning: pptx computed height is less than the minimum of 2,54cm. The slide heights will be set on 2,54cm.")

    #create slide
    prs = Presentation()
    prs.slide_height = Inches(calculate.mm_to_inch(figure_height))
    prs.slide_width = Inches(calculate.mm_to_inch(sum_total_width_mm))
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # place modules and their corresponding elements and save
    place_modules(data, to_path, slide)
    prs.save(filename)