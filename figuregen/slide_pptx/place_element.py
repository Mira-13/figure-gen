from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from . import calculate

class Error(Exception):
    def __init__(self, message):
        self.message = message

def add_image(slide, path, width_inch, pos_top, pos_left):
    slide.shapes.add_picture(path, Inches(pos_left), Inches(pos_top), width=Inches(width_inch))

def create_rectangle(slide, pos_top, pos_left, width, height):
    '''
    position and size in inches.
    Default: no shadow
    '''
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(pos_left), Inches(pos_top), Inches(width), Inches(height))
    shape.shadow.inherit = False
    return shape

def _apply_fill_color(shape, color):
    if color is None:
        shape.fill.background() # = no fill color
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(color[0], color[1], color[2])
    shape.line.fill.background()

def _add_frame_on_top(slide, pos_top, pos_left, width_inch, height_inch, color, thickness_pt):
    '''
    Frames are rectangles placed on top of an image element. 
    A frame has: 
        - no background color, 
        - no text
    '''
    t_inch = calculate.pt_to_inch(thickness_pt)
    shape = create_rectangle(slide, pos_top + t_inch/2. , pos_left + t_inch/2, width_inch - t_inch, height_inch - t_inch)

    shape.line.color.rgb = RGBColor(color[0], color[1], color[2])
    shape.line.width = Pt(thickness_pt)
    #shape.line.join_type = 'Miter' # Removes rounded edges, but is not supported, yet (sadly) 
    shape.fill.background()

def _apply_text_alignment(paragraph, alignment):
    if alignment == 'center' or alignment == 'centering':
        paragraph.alignment = PP_ALIGN.CENTER
    elif alignment == 'left':
        paragraph.alignment = PP_ALIGN.LEFT
    elif alignment == 'right':
        paragraph.alignment = PP_ALIGN.RIGHT
    else:
        raise Error('Could not apply text aligment. Incorrect value: '+alignment+'.')

def _apply_margins(text_frame, alignment, text_padding_mm):
    '''
    Margins of top and bottom are per default zero. Depending on the alignment, left or right margin will be set with the given text padding.
    '''
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    if alignment == 'right':
        text_frame.margin_right = Inches(calculate.mm_to_inch(text_padding_mm))
    else:
        text_frame.margin_right = 0
    if alignment == 'left':
        text_frame.margin_left = Inches(calculate.mm_to_inch(text_padding_mm))
    else:
        text_frame.margin_left = 0

def _apply_text_properties(shape, text, fontsize, txt_color, alignment, text_padding):
    '''
    Insert text into shape, set fontsize, and applies text color.
    '''
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    _apply_text_alignment(p, alignment)
    _apply_margins(text_frame, alignment, text_padding)

    run = p.add_run()
    run.text = text

    font = run.font
    font.color.rgb = RGBColor(txt_color[0], txt_color[1], txt_color[2])
    font.size = Pt(fontsize)

def add_text(slide, pos_top, pos_left, width_inch, height_inch, text, txt_rotation, fontsize_pt, txt_color=[0,0,0], bg_color=None, alignment='center', text_padding=0):
    '''
    Add a rectangle which includes:
        - text properties: text content, rotation (only 0° / +- 90°), text color
        - box properties: background color (but no frame)

    Careful: pptx package does not support text rotation, yet.
    However, we support the rotation (+-) 90° by doing some workarounds:
    - we switch height and width of the rectangle, and 
    - rotate the whole shape afterwards
    --> the shape, including it's content (text) is therefore rotated.
    Other rotation values are not supported, except for 0°. 
    '''
    if txt_rotation == 90.0 or txt_rotation == -90.0:
        # The shape is rotated about its center. We want a rotation about the top left corner instead.
        # Since we only allow 90° rotations, we can correct for that with a simple translation
        pos_top += height_inch / 2. - width_inch / 2.
        pos_left -= height_inch / 2. - width_inch / 2.
        
        # swap height and width
        height_inch, width_inch = width_inch, height_inch

        shape = create_rectangle(slide, pos_top, pos_left, width_inch, height_inch)
        shape.rotation = -txt_rotation # tikz rotation is counter-clockwise, pptx clockwise (we switch in pptx)
    else:
        shape = create_rectangle(slide, pos_top, pos_left, width_inch, height_inch)
    _apply_fill_color(shape, bg_color)
    _apply_text_properties(shape, text, fontsize_pt, txt_color, alignment, text_padding)

def _add_label(slide, img_pos_top, img_pos_left, img_width_inch, img_height_inch, cfg, label_pos, factor):
    if cfg is None:
        return
    try:
        cfg = cfg[label_pos]
    except KeyError:
        return
                
    alignment = label_pos.split('_')[-1]
    is_top = (label_pos.split('_')[0] == 'top')
       
    rect_width, rect_height = calculate.mm_to_inch(cfg['width_mm'])*factor, calculate.mm_to_inch(cfg['height_mm'])*factor

    # determine the correct offsets depending on wether it is in the corner or center
    if alignment == 'center':
        offset_w, offset_h = 0, calculate.mm_to_inch(cfg['offset_mm'])*factor
    else:
        offset_w = calculate.mm_to_inch(cfg['offset_mm'][0])*factor
        offset_h = calculate.mm_to_inch(cfg['offset_mm'][1])*factor

    # determine pos_top and pos_left of rectangle
    if is_top:
        pos_top = img_pos_top + offset_h
    else:
        pos_top = img_pos_top + img_height_inch - rect_height - offset_h
    if alignment == 'center':
        pos_left = img_pos_left + (img_width_inch * 1/2) - (rect_width * 1/2)
    elif alignment == 'left':
        pos_left = img_pos_left + offset_w
    else: # right
        pos_left = img_pos_left + img_width_inch - rect_width - offset_w

    add_text(slide, pos_top, pos_left, rect_width, rect_height, cfg['text'], 0, cfg['fontsize'], cfg['text_color'], cfg['background_color'], alignment, cfg['padding_mm'])

def _add_labels(slide, img_pos_top, img_pos_left, img_width_inch, img_height_inch, cfg, factor):
    for label_pos in ['top_center', 'top_left', 'top_right', 'bottom_center', 'bottom_left', 'bottom_right']:
        _add_label(slide, img_pos_top, img_pos_left, img_width_inch, img_height_inch, cfg, label_pos, factor)

def _has_frame(element):
    try:
        frame = element['frame']
    except:
        return False
    return True

def _add_markers(slide, element, img_pos_top, img_pos_left, img_width_px, img_height_px, img_used_width, img_used_height):
    try:
        markers = element['marker']
    except:
        return

    # crop_markers are based on pixels, therefore we calculate the relative position of the marker that will be placed on top of the image
    w_scale, h_scale = calculate.relative_position(img_width_px, img_height_px, img_used_width, img_used_height)

    if markers['list']!=[]: # only draw if line width reasonable and list not empty
        for m in markers['list']:
            if m['lw'] > 0.0:
                pos_top = img_pos_top + m['pos'][1] * h_scale
                pos_left = img_pos_left + m['pos'][0] * w_scale
                w = m['size'][0] * w_scale
                h = m['size'][1] * h_scale
                _add_frame_on_top(slide, pos_top, pos_left, w, h, m['color'], m['lw'])

def _add_lines(slide, element, img_pos_top, img_pos_left, img_width_px, img_height_px, img_used_width, img_used_height):
    try:
        lines = element['lines']
    except:
        return

    w_scale, h_scale = calculate.relative_position(img_width_px, img_height_px, img_used_width, img_used_height)

    if lines != []: # only draw if line width reasonable and list not empty
        for l in lines:
            if l['lw'] > 0.0:
                pos_start_h = img_pos_top + l['from'][0] * h_scale
                pos_start_w = img_pos_left + l['from'][1] * w_scale
                pos_end_h = img_pos_top + l['to'][0] * h_scale
                pos_end_w = img_pos_left + l['to'][1] * w_scale
                shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(pos_start_w), Inches(pos_start_h), Inches(pos_end_w), Inches(pos_end_h))
                shape.shadow.inherit = False
                shape.line.color.rgb = RGBColor(l['color'][0], l['color'][1], l['color'][2])
                shape.line.width = Pt(l['lw'])

def images_and_frames_and_labels(slide, data, factor, offset_width_mm, offset_top_mm):
    '''
    Reads module data and puts images on the slide. 
    args:
        factor; is slide_width / figure_width, which is used to position elements on the slide accordingly
    '''
    width_inch, height_inch = calculate.img_size_inches(data, factor)

    rowIndex = 1
    for row in data['elements_content']:
        colIndex = 1
        for element in row:
            # place image
            pos_top, pos_left = calculate.img_pos(data, colIndex, rowIndex, factor, offset_left_mm=offset_width_mm, offset_top_mm=offset_top_mm)
            add_image(slide, element['filename'], width_inch, pos_top, pos_left)

            # place frame
            if _has_frame(element):
                _add_frame_on_top(slide, pos_top, pos_left, width_inch, height_inch, color=element['frame']['color'], thickness_pt=element['frame']['line_width'])
            
            # place labels
            try:
                cfg = element['label']
            except:
                cfg = None
            _add_labels(slide, pos_top, pos_left, width_inch, height_inch, cfg, factor)

            # place markers
            _add_markers(slide, element, pos_top, pos_left, data['img_width_px'], data['img_height_px'], width_inch, height_inch)
            _add_lines(slide, element, pos_top, pos_left, data['img_width_px'], data['img_height_px'], width_inch, height_inch)

            colIndex += 1
        rowIndex += 1


def titles(slide, data, factor, offset_left_mm, offset_top_mm=0.0):
    for direction in ['north', 'east', 'south', 'west']:
        position, size = calculate.titles_pos(data, direction, factor, offset_left_mm, offset_top_mm)
        if size[0] != 0.0 and size[1] != 0.0:
            title = data['titles'][direction]
            add_text(slide, position[0], position[1], size[0], size[1], 
                     title['content'], title['rotation'], title['fontsize'], 
                     title['text_color'], title['background_color'])

def _compute_bg_colors(bg_color_properties, num):
    if bg_color_properties is None: # no background color
        return [None for i in range(num)]
    elif not isinstance(bg_color_properties[0], list): # constant color for all
        return [bg_color_properties for i in range(num)]
    else: # individual background colors
        return bg_color_properties

def _row_col_titles(slide, data, direction, title_properties, num, pos_fn):
    if calculate.size_of(title_properties, direction)[0] != 0.0:
        bg_colors = _compute_bg_colors(title_properties[direction]['background_colors'], num)
        t = title_properties[direction]
        for i in range(num):
            position, size = pos_fn(i)
            add_text(slide, *position, *size, t['content'][i], t['rotation'], t['fontsize'], t['text_color'], bg_colors[i])
  
def row_titles(slide, data, factor, offset_left_mm, offset_top_mm=0.0):
    for direction in ['east', 'west']:
        def pos_fn (idx): 
            return calculate.row_titles_pos(data, idx + 1, direction, factor, offset_left_mm, offset_top_mm)
        _row_col_titles(slide, data, direction, data['row_titles'], data['num_rows'], pos_fn)

def col_titles(slide, data, factor, offset_left_mm, offset_top_mm=0.0):
    for direction in ['north', 'south']:
        def pos_fn (idx): 
            return calculate.column_titles_pos(data, idx + 1, direction, factor, offset_left_mm, offset_top_mm)
        _row_col_titles(slide, data, direction, data['column_titles'], data['num_columns'], pos_fn)

def south_captions(slide, data, factor, offset_left_mm, offset_top_mm=0.0):
    '''
    #new feature: allow south captions (below) each image.
    For now, we ignore the background-color for those captions.
    '''
    capt_prop = data['element_config']['captions']['south']
    size = calculate.img_size_inches(data, factor)[0], calculate.mm_to_inch(capt_prop['height']) * factor
    if size[0] == 0.0:
        return
    
    rowIndex = 1
    for row in data['elements_content']:
        colIndex = 1
        for elem in row:
            position = calculate.south_caption_pos(data, colIndex, rowIndex, factor, offset_left_mm, offset_top_mm)
            txt_content = elem['captions']['south']
            add_text(slide, *position, *size, txt_content, capt_prop['rotation'], capt_prop['fontsize'], capt_prop['text_color']) 
            colIndex += 1
        rowIndex += 1